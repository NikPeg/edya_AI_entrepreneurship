#!/usr/bin/env python3
"""
Универсальный генератор презентаций из YAML конфига
"""
import sys
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, config):
    """Слайд 1: Титульный"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = config['title']
    subtitle.text = f"{config['subtitle']}\n{', '.join([m['name'] for m in config['team']])}\n{config['date']}"

def add_team_slide(prs, config):
    """Слайд 2: Команда"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "КОМАНДА"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for member in config['team']:
        p = tf.add_paragraph()
        p.text = f"{member['name']}"
        p.level = 0
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = f"{member['role']}"
        p.level = 1

def add_idea_slides(prs, config):
    """Слайд 3+: Идеи продукта (1 идея = 1 слайд)"""
    for idea in config['ideas']:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "ИДЕЯ ПРОДУКТА"
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        # Название идеи
        p = tf.add_paragraph()
        p.text = idea['title']
        p.font.bold = True
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = ""
        
        # UVP (если есть)
        if 'uvp' in idea:
            # Разбиваем по пустым строкам (абзацам), а не по каждой строке
            paragraphs = idea['uvp'].strip().split('\n\n')
            for paragraph in paragraphs:
                if paragraph.strip():
                    # Объединяем строки внутри абзаца
                    text = ' '.join(line.strip() for line in paragraph.split('\n') if line.strip())
                    p = tf.add_paragraph()
                    p.text = text
                    p.level = 0
                    p.font.size = Pt(12)
        # Старый формат с description (для обратной совместимости)
        elif 'description' in idea:
            paragraphs = idea['description'].strip().split('\n\n')
            for paragraph in paragraphs:
                if paragraph.strip():
                    text = ' '.join(line.strip() for line in paragraph.split('\n') if line.strip())
                    p = tf.add_paragraph()
                    p.text = text
                    p.level = 0
        
        # Источники
        if 'sources' in idea:
            p = tf.add_paragraph()
            p.text = ""
            p = tf.add_paragraph()
            p.text = "Источники:"
            p.font.bold = True
            p.font.size = Pt(11)
            
            for source in idea['sources']:
                p = tf.add_paragraph()
                p.text = source
                p.level = 1
                p.font.size = Pt(9)

def add_icp_slide(prs, config):
    """Слайд 4: ICP - Профиль клиента"""
    icp = config['icp']
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "LEAN CANVAS. ICP - ГИПОТЕЗА ПРОФИЛЯ КЛИЕНТА"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    # Сегмент
    p = tf.add_paragraph()
    p.text = f"Сегмент: {icp['segment']}"
    p.font.bold = True
    
    # Размер компании
    p = tf.add_paragraph()
    p.text = f"Размер: {icp['company_size']}"
    
    # Отрасли
    p = tf.add_paragraph()
    p.text = "Отрасли:"
    p.font.bold = True
    for industry in icp['industries']:
        p = tf.add_paragraph()
        p.text = industry
        p.level = 1
    
    # Боли клиентов
    p = tf.add_paragraph()
    p.text = "Основные боли:"
    p.font.bold = True
    for pain in icp['pain_points'][:3]:  # Первые 3 боли
        p = tf.add_paragraph()
        p.text = pain
        p.level = 1
    
    # Decision makers
    p = tf.add_paragraph()
    p.text = f"Лица принимающие решение: {', '.join(icp['decision_makers'])}"
    p.font.size = Pt(12)

def add_problem_slide(prs, config):
    """Слайд 5: Проблема"""
    problem = config['lean_canvas']['problem']
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"LEAN CANVAS. {problem['title'].upper()}"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in problem['items']:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

def add_solution_slide(prs, config):
    """Слайд 6: Решение"""
    solution = config['lean_canvas']['solution']
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"LEAN CANVAS. {solution['title'].upper()}"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in solution['items']:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

def add_uvp_slide(prs, config):
    """Слайд 7: UVP"""
    uvp = config['lean_canvas']['uvp']
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"LEAN CANVAS. {uvp['title'].upper()}"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    # Основное УТП
    p = tf.add_paragraph()
    p.text = uvp['main']
    p.font.bold = True
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = ""
    
    # Детали
    for detail in uvp['details']:
        p = tf.add_paragraph()
        p.text = detail
        p.level = 0

def add_channels_slide(prs, config):
    """Слайд 8: Каналы"""
    channels = config['lean_canvas']['channels']
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"LEAN CANVAS. {channels['title'].upper()}"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in channels['items']:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

def add_revenue_cost_slide(prs, config):
    """Слайд 9: Доходы и расходы"""
    revenue = config['lean_canvas']['revenue_streams']
    cost = config['lean_canvas']['cost_structure']
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "LEAN CANVAS. СТРУКТУРА ДОХОДОВ И РАСХОДОВ"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    # Доходы
    p = tf.add_paragraph()
    p.text = revenue['title'].upper()
    p.font.bold = True
    p.font.size = Pt(14)
    
    for item in revenue['items']:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(11)
    
    if 'target' in revenue:
        p = tf.add_paragraph()
        p.text = f"Цель: {revenue['target']}"
        p.font.bold = True
        p.font.size = Pt(11)
    
    p = tf.add_paragraph()
    p.text = ""
    
    # Расходы
    p = tf.add_paragraph()
    p.text = cost['title'].upper()
    p.font.bold = True
    p.font.size = Pt(14)
    
    for item in cost['items'][:4]:  # Первые 4
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(11)

def add_metrics_slide(prs, config):
    """Слайд 10: Ключевые метрики"""
    metrics = config['lean_canvas']['key_metrics']
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"LEAN CANVAS. {metrics['title'].upper()}"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in metrics['items']:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

def add_unfair_advantage_slide(prs, config):
    """Слайд 11: Нечестное преимущество"""
    advantage = config['lean_canvas']['unfair_advantage']
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"LEAN CANVAS. {advantage['title'].upper()}"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in advantage['items']:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

def add_segments_slide(prs, config):
    """Дополнительный слайд: Сегменты клиентов (если нужен)"""
    segments = config['lean_canvas']['customer_segments']
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"LEAN CANVAS. {segments['title'].upper()}"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    # Ранние последователи
    p = tf.add_paragraph()
    p.text = "Ранние последователи:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = segments['early_adopters']
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    
    # Сегменты
    p = tf.add_paragraph()
    p.text = "Сегменты рынка:"
    p.font.bold = True
    
    for segment in segments['segments']:
        p = tf.add_paragraph()
        p.text = segment
        p.level = 0

def generate_presentation(config_path, output_path):
    """Главная функция генерации"""
    print(f"Читаю конфиг: {config_path}")
    
    with open(config_path, 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    print("Создаю презентацию...")
    prs = Presentation()
    
    # Генерируем слайды по порядку
    add_title_slide(prs, config)
    add_team_slide(prs, config)
    add_idea_slides(prs, config)
    add_icp_slide(prs, config)
    add_problem_slide(prs, config)
    add_solution_slide(prs, config)
    add_uvp_slide(prs, config)
    add_channels_slide(prs, config)
    add_revenue_cost_slide(prs, config)
    add_metrics_slide(prs, config)
    add_unfair_advantage_slide(prs, config)
    add_segments_slide(prs, config)
    
    # Сохраняем
    prs.save(output_path)
    print(f"✅ Презентация сохранена: {output_path}")
    print(f"   Слайдов создано: {len(prs.slides)}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Использование: python generate_pptx.py <config.yaml> [output.pptx]")
        print("\nПример:")
        print("  python scripts/generate_pptx.py homework/homework_1/presentation_config.yaml")
        print("  python scripts/generate_pptx.py config.yaml output/presentation.pptx")
        sys.exit(1)
    
    config_path = sys.argv[1]
    
    if len(sys.argv) > 2:
        output_path = sys.argv[2]
    else:
        # Автоматическое имя: config.yaml -> config.pptx
        output_path = config_path.replace('.yaml', '.pptx').replace('.yml', '.pptx')
    
    generate_presentation(config_path, output_path)

