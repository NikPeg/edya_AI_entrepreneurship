#!/usr/bin/env python3
"""
Скрипт для извлечения текста из PPTX файлов
"""
import sys
from pptx import Presentation

def pptx_to_txt(pptx_path, txt_path=None):
    """
    Извлекает текст из PPTX файла и сохраняет в TXT
    
    Args:
        pptx_path: путь к PPTX файлу
        txt_path: путь к выходному TXT файлу (если None, будет создан автоматически)
    """
    if txt_path is None:
        txt_path = pptx_path.replace('.pptx', '.txt')
    
    print(f"Читаю PPTX: {pptx_path}")
    
    try:
        # Открываем презентацию
        prs = Presentation(pptx_path)
        
        print(f"Найдено слайдов: {len(prs.slides)}")
        
        with open(txt_path, 'w', encoding='utf-8') as txt_file:
            for i, slide in enumerate(prs.slides, 1):
                print(f"Обрабатываю слайд {i}/{len(prs.slides)}", end='\r')
                
                # Записываем заголовок слайда
                txt_file.write(f"\n{'='*80}\n")
                txt_file.write(f"СЛАЙД {i}\n")
                txt_file.write(f"{'='*80}\n\n")
                
                # Извлекаем текст из всех shapes на слайде
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            txt_file.write(text + "\n")
                
                txt_file.write("\n")
        
        print(f"\n\nГотово! Текст сохранен в: {txt_path}")
        return txt_path
        
    except Exception as e:
        print(f"\nОшибка при обработке PPTX: {e}")
        print("\nУбедитесь, что установлена библиотека python-pptx:")
        print("  pip install python-pptx")
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Использование: python pptx_to_txt.py <путь_к_pptx> [путь_к_txt]")
        print("\nПример:")
        print("  python pptx_to_txt.py presentation.pptx")
        print("  python pptx_to_txt.py presentation.pptx output.txt")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    txt_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    pptx_to_txt(pptx_file, txt_file)

